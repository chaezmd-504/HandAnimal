from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Color palette ──────────────────────────────────────────────
BG      = RGBColor(0x0D, 0x1B, 0x2A)   # dark navy
ACCENT  = RGBColor(0x00, 0xC8, 0xFF)   # cyan
ACCENT2 = RGBColor(0x7B, 0xFF, 0xD4)   # mint
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0xAA, 0xBB, 0xCC)
YELLOW  = RGBColor(0xFF, 0xD7, 0x00)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank

# ── Helpers ────────────────────────────────────────────────────
def add_slide():
    slide = prs.slides.add_slide(blank_layout)
    # dark background rectangle
    bg = slide.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    return slide

def txbox(slide, text, x, y, w, h,
          size=24, bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return tb

def accent_bar(slide, x, y, w=Inches(0.06), h=Inches(0.7)):
    bar = slide.shapes.add_shape(1, x, y, w, h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

def divider(slide, y):
    line = slide.shapes.add_shape(1, Inches(0.6), y, Inches(12.1), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

def bullet_box(slide, items, x, y, w, h, size=19, color=WHITE, gap=Pt(6)):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = gap
        run = p.add_run()
        run.text = item
        run.font.size  = Pt(size)
        run.font.color.rgb = color
    return tb

def tag(slide, text, x, y, w=Inches(2.8), h=Inches(0.5)):
    box = slide.shapes.add_shape(1, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x00, 0x50, 0x80)
    box.line.color.rgb = ACCENT
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = False
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = ACCENT

def slide_num(slide, n, total=9):
    txbox(slide, f"{n} / {total}",
          Inches(11.8), Inches(7.1), Inches(1.4), Inches(0.4),
          size=12, color=GRAY, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════
# Slide 1 — Title
# ══════════════════════════════════════════════════════════════════
s = add_slide()

# gradient-like accent strip on left
strip = s.shapes.add_shape(1, 0, 0, Inches(0.35), H)
strip.fill.solid(); strip.fill.fore_color.rgb = ACCENT
strip.line.fill.background()

# paw-print placeholder circle (decorative)
for cx, cy, r in [(Inches(11.5), Inches(1.2), Inches(0.5)),
                   (Inches(12.3), Inches(2.1), Inches(0.35)),
                   (Inches(11.0), Inches(2.8), Inches(0.28))]:
    c = s.shapes.add_shape(9, cx, cy, r, r)   # oval
    c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0x00, 0x40, 0x60)
    c.line.fill.background()

txbox(s, "HandAnimal",
      Inches(0.7), Inches(1.6), Inches(10), Inches(1.6),
      size=60, bold=True, color=ACCENT)

txbox(s, "Real-Time Animal Control via Hand & Gaze Tracking",
      Inches(0.7), Inches(3.1), Inches(10.5), Inches(1.0),
      size=28, bold=False, color=WHITE)

divider(s, Inches(4.2))

txbox(s, "Intuitive embodied control of virtual animals\nthrough natural hand and gaze movement",
      Inches(0.7), Inches(4.4), Inches(10.5), Inches(1.2),
      size=20, color=GRAY)

slide_num(s, 1)

# ══════════════════════════════════════════════════════════════════
# Slide 2 — Problem & Motivation
# ══════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, Inches(0.6), Inches(0.6))
txbox(s, "Problem & Motivation",
      Inches(0.85), Inches(0.55), Inches(9), Inches(0.7),
      size=32, bold=True, color=ACCENT)
divider(s, Inches(1.4))

# Left column — problem
txbox(s, "Traditional Approach",
      Inches(0.6), Inches(1.6), Inches(5.5), Inches(0.5),
      size=20, bold=True, color=YELLOW)
bullet_box(s, [
    "⌨  Keyboard / Controller input",
    "⚠  No sense of body ownership",
    "⚠  Breaks immersion",
    "⚠  High learning curve",
], Inches(0.7), Inches(2.1), Inches(5.2), Inches(2.5), size=19)

# Right column — solution
txbox(s, "HandAnimal Approach",
      Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.5),
      size=20, bold=True, color=ACCENT2)
bullet_box(s, [
    "🖐  Natural hand gestures",
    "👁  Gaze-based navigation",
    "✅  Embodied body ownership",
    "✅  Zero learning curve",
], Inches(7.0), Inches(2.1), Inches(5.8), Inches(2.5), size=19)

# Arrow in middle
txbox(s, "→", Inches(6.0), Inches(2.6), Inches(1.0), Inches(0.8),
      size=48, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

txbox(s, '"What if you could wear the animal\'s body?"',
      Inches(1.5), Inches(5.5), Inches(10), Inches(0.8),
      size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

slide_num(s, 2)

# ══════════════════════════════════════════════════════════════════
# Slide 3 — System Overview
# ══════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, Inches(0.6), Inches(0.6))
txbox(s, "System Overview",
      Inches(0.85), Inches(0.55), Inches(9), Inches(0.7),
      size=32, bold=True, color=ACCENT)
divider(s, Inches(1.4))

# Pipeline nodes
nodes = [
    ("Camera", Inches(0.5)),
    ("Python\n(MediaPipe)", Inches(2.5)),
    ("WebSocket", Inches(4.9)),
    ("Unity\nEngine", Inches(7.2)),
    ("Animal\nController", Inches(9.7)),
]
ny = Inches(2.4); nw = Inches(1.7); nh = Inches(1.1)
for label, nx in nodes:
    box = s.shapes.add_shape(1, nx, ny, nw, nh)
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x00, 0x4A, 0x7C)
    box.line.color.rgb = ACCENT
    tb = s.shapes.add_textbox(nx, ny, nw, nh)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = label
    run.font.size = Pt(16); run.font.bold = True; run.font.color.rgb = WHITE

# Arrows between nodes
for i in range(len(nodes) - 1):
    ax = nodes[i][1] + nw
    txbox(s, "→", ax, ny, Inches(0.6), nh,
          size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# Two channel labels
txbox(s, "Channel 1 — Hand Tracking",
      Inches(0.5), Inches(3.8), Inches(5.5), Inches(0.5),
      size=16, color=ACCENT2)
txbox(s, "Channel 2 — Gaze Tracking",
      Inches(0.5), Inches(4.2), Inches(5.5), Inches(0.5),
      size=16, color=YELLOW)

# Tech stack tags
tags = ["MediaPipe Hands", "MediaPipe FaceMesh", "Python 3.10",
        "WebSockets", "Unity 2022", "C# Scripts"]
for i, t in enumerate(tags):
    tx = Inches(0.5 + (i % 3) * 4.2)
    ty = Inches(5.3 + (i // 3) * 0.65)
    tag(s, t, tx, ty, Inches(3.8), Inches(0.5))

slide_num(s, 3)

# ══════════════════════════════════════════════════════════════════
# Slide 4 — Hand Tracking & Joint Mapping
# ══════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, Inches(0.6), Inches(0.6))
txbox(s, "Hand Tracking & Joint Mapping",
      Inches(0.85), Inches(0.55), Inches(11), Inches(0.7),
      size=32, bold=True, color=ACCENT)
divider(s, Inches(1.4))

# Left — concept
txbox(s, "Input → Animal Joint",
      Inches(0.6), Inches(1.6), Inches(5.5), Inches(0.5),
      size=20, bold=True, color=ACCENT2)
mappings = [
    "Finger curl   →   Leg joint flexion",
    "Wrist angle   →   Spine / shoulder",
    "Palm normal   →   Body tilt",
    "Hand open     →   Rest / idle pose",
]
bullet_box(s, mappings, Inches(0.7), Inches(2.1), Inches(5.5), Inches(2.5), size=18)

# Right — techniques
txbox(s, "Key Techniques",
      Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.5),
      size=20, bold=True, color=YELLOW)
techniques = [
    "📐  ROM-based joint mapping",
    "〰  EMA smoothing (2-stage)",
    "🔀  Blend mode: idle / walk / action",
    "⚖  Left-right ROM symmetrization",
    "🎯  5° dead-zone threshold",
]
bullet_box(s, techniques, Inches(7.0), Inches(2.1), Inches(5.8), Inches(3.0), size=18)

# Supported animals
txbox(s, "Supported Animals",
      Inches(0.6), Inches(5.0), Inches(5), Inches(0.5),
      size=18, bold=True, color=GRAY)
for i, animal in enumerate(["🕷  Spider  (8-leg IK)", "🐴  Horse  (quadruped)"]):
    txbox(s, animal, Inches(0.7 + i * 5.5), Inches(5.5), Inches(5), Inches(0.5),
          size=20, bold=True, color=WHITE)

slide_num(s, 4)

# ══════════════════════════════════════════════════════════════════
# Slide 5 — Locomotion Control
# ══════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, Inches(0.6), Inches(0.6))
txbox(s, "Locomotion Control",
      Inches(0.85), Inches(0.55), Inches(9), Inches(0.7),
      size=32, bold=True, color=ACCENT)
divider(s, Inches(1.4))

# Flow
flow = [
    ("Left Hand\nVelocity", RGBColor(0x00,0x4A,0x7C)),
    ("Speed\nMapping", RGBColor(0x00,0x5A,0x50)),
    ("Animation\nBlend", RGBColor(0x3A,0x3A,0x00)),
    ("Animal\nMovement", RGBColor(0x50,0x00,0x50)),
]
fy = Inches(2.0); fw = Inches(2.2); fh = Inches(1.2)
for i, (label, color) in enumerate(flow):
    fx = Inches(0.5 + i * 3.1)
    box = s.shapes.add_shape(1, fx, fy, fw, fh)
    box.fill.solid(); box.fill.fore_color.rgb = color
    box.line.color.rgb = ACCENT
    tb = s.shapes.add_textbox(fx, fy, fw, fh)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = label
    run.font.size = Pt(17); run.font.bold = True; run.font.color.rgb = WHITE
    if i < len(flow) - 1:
        txbox(s, "→", fx + fw, fy, Inches(0.9), fh,
              size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# Spider details
txbox(s, "Spider", Inches(0.6), Inches(3.7), Inches(5.5), Inches(0.5),
      size=20, bold=True, color=ACCENT2)
bullet_box(s, [
    "Walk animation speed synced to hand velocity",
    "Smooth transition: idle ↔ walk ↔ run",
    "8 legs animated via procedural IK",
], Inches(0.7), Inches(4.2), Inches(5.5), Inches(2.0), size=18)

# Horse details
txbox(s, "Horse", Inches(7.0), Inches(3.7), Inches(5.5), Inches(0.5),
      size=20, bold=True, color=YELLOW)
bullet_box(s, [
    "Side override for lateral movement",
    "Gait blend: walk / trot / canter",
    "Head bow & body lean response",
], Inches(7.0), Inches(4.2), Inches(5.8), Inches(2.0), size=18)

slide_num(s, 5)

# ══════════════════════════════════════════════════════════════════
# Slide 6 — Gaze-Based Navigation
# ══════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, Inches(0.6), Inches(0.6))
txbox(s, "Gaze-Based Navigation",
      Inches(0.85), Inches(0.55), Inches(9), Inches(0.7),
      size=32, bold=True, color=ACCENT)
divider(s, Inches(1.4))

# Three component boxes
comps = [
    ("MediaPipe\nFaceMesh", "Detects iris position\n& gaze direction", ACCENT),
    ("GazeCalibrator", "User-specific\ncalibration\n(min/max gaze range)", ACCENT2),
    ("GazeNavigator", "Maps gaze vector\nto animal heading\nwith smooth blend", YELLOW),
]
cx_list = [Inches(0.5), Inches(4.6), Inches(8.7)]
for (title, desc, color), cx in zip(comps, cx_list):
    cw = Inches(3.8); ch = Inches(2.2); cy = Inches(2.0)
    box = s.shapes.add_shape(1, cx, cy, cw, ch)
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x00,0x2A,0x40)
    box.line.color.rgb = color
    txbox(s, title, cx + Inches(0.1), cy + Inches(0.1), cw, Inches(0.7),
          size=18, bold=True, color=color)
    txbox(s, desc, cx + Inches(0.1), cy + Inches(0.8), cw - Inches(0.2), Inches(1.3),
          size=16, color=WHITE)
    if cx != cx_list[-1]:
        txbox(s, "→", cx + cw, cy + Inches(0.7), Inches(0.7), Inches(0.7),
              size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

txbox(s, "Data flow",
      Inches(0.5), Inches(4.5), Inches(3), Inches(0.4),
      size=16, bold=True, color=GRAY)
bullet_box(s, [
    "Eye landmark → normalized gaze (x, y)",
    "Calibrated range → direction vector",
    "EMA smoothing → stable heading output",
    "WebSocket → Unity GazeNavigator in real-time",
], Inches(0.6), Inches(4.9), Inches(12), Inches(2.0), size=18)

slide_num(s, 6)

# ══════════════════════════════════════════════════════════════════
# Slide 7 — Demo & Results
# ══════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, Inches(0.6), Inches(0.6))
txbox(s, "Demo & Results",
      Inches(0.85), Inches(0.55), Inches(9), Inches(0.7),
      size=32, bold=True, color=ACCENT)
divider(s, Inches(1.4))

# Two demo panels
for i, (animal, color) in enumerate([("🕷  Spider Control", ACCENT2), ("🐴  Horse Control", YELLOW)]):
    px = Inches(0.5 + i * 6.6); py = Inches(1.7)
    pw = Inches(6.0); ph = Inches(4.2)
    panel = s.shapes.add_shape(1, px, py, pw, ph)
    panel.fill.solid(); panel.fill.fore_color.rgb = RGBColor(0x05,0x10,0x20)
    panel.line.color.rgb = color

    txbox(s, animal, px + Inches(0.1), py + Inches(0.15), pw, Inches(0.6),
          size=22, bold=True, color=color)

    # placeholder image area
    img_box = s.shapes.add_shape(1, px + Inches(0.2), py + Inches(0.75),
                                  pw - Inches(0.4), Inches(2.0))
    img_box.fill.solid(); img_box.fill.fore_color.rgb = RGBColor(0x10,0x20,0x30)
    img_box.line.color.rgb = color

    txbox(s, "[ Demo Screenshot ]",
          px + Inches(0.2), py + Inches(1.1), pw - Inches(0.4), Inches(0.6),
          size=14, color=GRAY, align=PP_ALIGN.CENTER)

    feats = (["Blend mode switching", "8-leg IK locomotion", "Gaze steering"]
             if i == 0 else
             ["Gait blend walk/trot", "Head bow & body lean", "Side override"])
    bullet_box(s, feats, px + Inches(0.2), py + Inches(2.85),
               pw - Inches(0.4), Inches(1.3), size=16)

slide_num(s, 7)

# ══════════════════════════════════════════════════════════════════
# Slide 8 — Challenges & Solutions
# ══════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, Inches(0.6), Inches(0.6))
txbox(s, "Challenges & Solutions",
      Inches(0.85), Inches(0.55), Inches(9), Inches(0.7),
      size=32, bold=True, color=ACCENT)
divider(s, Inches(1.4))

challenges = [
    ("Joint Noise",       "Jittery landmark detection\ncauses uncontrolled movement",
     "Per-joint 2-stage EMA smoothing\nindependent alpha per joint"),
    ("L/R Asymmetry",     "Hand skeleton is not\nperfectly symmetric",
     "ROM symmetrization:\nmirror mapping calibration"),
    ("Latency",           "Python → Unity round-trip\naffects real-time feel",
     "WebSocket low-latency stream\n--fast / --skip low-spec flags"),
    ("Dead-zone Tuning",  "Micro-movements trigger\nunintended animation",
     "5° angular dead-zone\nfor stable idle state"),
]
for i, (title, prob, sol) in enumerate(challenges):
    ry = Inches(1.65 + i * 1.35); rh = Inches(1.2)
    # problem side
    pb = s.shapes.add_shape(1, Inches(0.5), ry, Inches(4.5), rh)
    pb.fill.solid(); pb.fill.fore_color.rgb = RGBColor(0x30,0x08,0x08)
    pb.line.color.rgb = RGBColor(0xFF,0x44,0x44)
    txbox(s, f"⚠  {title}", Inches(0.6), ry + Inches(0.05), Inches(4.3), Inches(0.45),
          size=16, bold=True, color=RGBColor(0xFF,0x88,0x88))
    txbox(s, prob, Inches(0.6), ry + Inches(0.48), Inches(4.3), Inches(0.65),
          size=14, color=WHITE)
    # arrow
    txbox(s, "→", Inches(5.1), ry + Inches(0.3), Inches(0.8), Inches(0.6),
          size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    # solution side
    sb = s.shapes.add_shape(1, Inches(6.0), ry, Inches(6.8), rh)
    sb.fill.solid(); sb.fill.fore_color.rgb = RGBColor(0x00,0x30,0x20)
    sb.line.color.rgb = ACCENT2
    txbox(s, f"✅  Solution", Inches(6.1), ry + Inches(0.05), Inches(6.5), Inches(0.45),
          size=16, bold=True, color=ACCENT2)
    txbox(s, sol, Inches(6.1), ry + Inches(0.48), Inches(6.5), Inches(0.65),
          size=14, color=WHITE)

slide_num(s, 8)

# ══════════════════════════════════════════════════════════════════
# Slide 9 — Conclusion & Future Work
# ══════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, Inches(0.6), Inches(0.6))
txbox(s, "Conclusion & Future Work",
      Inches(0.85), Inches(0.55), Inches(9), Inches(0.7),
      size=32, bold=True, color=ACCENT)
divider(s, Inches(1.4))

# Summary quote box
qbox = s.shapes.add_shape(1, Inches(0.6), Inches(1.65), Inches(12.1), Inches(1.5))
qbox.fill.solid(); qbox.fill.fore_color.rgb = RGBColor(0x00,0x30,0x55)
qbox.line.color.rgb = ACCENT
txbox(s, '"HandAnimal enables intuitive embodied control of virtual animals\nthrough natural hand and gaze movement."',
      Inches(0.8), Inches(1.75), Inches(11.8), Inches(1.3),
      size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# What we achieved
txbox(s, "What We Achieved",
      Inches(0.6), Inches(3.4), Inches(5.5), Inches(0.5),
      size=20, bold=True, color=ACCENT2)
bullet_box(s, [
    "✅  Real-time hand → joint mapping pipeline",
    "✅  Two animal types: Spider & Horse",
    "✅  Gaze-based navigation system",
    "✅  Smooth locomotion & blend mode control",
], Inches(0.7), Inches(3.9), Inches(5.5), Inches(2.5), size=18)

# Future work
txbox(s, "Future Work",
      Inches(7.0), Inches(3.4), Inches(5.5), Inches(0.5),
      size=20, bold=True, color=YELLOW)
bullet_box(s, [
    "🔮  More animal types (bird, fish, insect)",
    "🔮  Full-body pose tracking",
    "🔮  VR / XR headset integration",
    "🔮  Multiplayer embodiment",
], Inches(7.0), Inches(3.9), Inches(5.8), Inches(2.5), size=18)

slide_num(s, 9)

# ── Save ───────────────────────────────────────────────────────
out = "HandAnimal_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
