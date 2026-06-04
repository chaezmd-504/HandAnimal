"""
make_progress_pptx.py
HandAvatar 프로젝트 진행 보고 발표자료 생성 (2026-06)
실행: python python/scripts/make_progress_pptx.py
출력: progress_report.pptx (루트 폴더)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── 색상 팔레트 ──────────────────────────────────────
BG_DARK    = RGBColor(0x12, 0x12, 0x1E)
ACCENT1    = RGBColor(0x5B, 0x9B, 0xD5)   # 파랑
ACCENT2    = RGBColor(0x70, 0xAD, 0x47)   # 초록
ACCENT3    = RGBColor(0xFF, 0xC0, 0x00)   # 노랑
ACCENT4    = RGBColor(0xFF, 0x66, 0x44)   # 주황빨강 (미해결/경고)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_GRAY  = RGBColor(0x2D, 0x2D, 0x3F)
SUBTITLE   = RGBColor(0xA0, 0xC4, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide(accent=ACCENT1):
    slide = prs.slides.add_slide(BLANK)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_DARK
    # 상단 바
    r = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.08))
    r.fill.solid(); r.fill.fore_color.rgb = accent; r.line.fill.background()
    return slide


def txb(slide, text, x, y, w, h,
        size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Malgun Gothic"
    return box


def add_para(txbox, text, size=15, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, space_before=Pt(4)):
    tf = txbox.text_frame
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Malgun Gothic"
    return p


def rect(slide, x, y, w, h, fill=DARK_GRAY, line=None, lw=Pt(1.5)):
    shape = slide.shapes.add_shape(1,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = lw
    else:
        shape.line.fill.background()
    return shape


def slide_title(slide, text, sub=None, accent=ACCENT1):
    txb(slide, text, 0.4, 0.15, 12.5, 0.75, size=34, bold=True, color=WHITE)
    if sub:
        txb(slide, sub, 0.4, 0.85, 12.5, 0.45, size=16, color=SUBTITLE)


def video_placeholder(slide, x, y, w, h, label="📹 영상 삽입"):
    rect(slide, x, y, w, h,
         fill=RGBColor(0x1A, 0x1A, 0x30),
         line=RGBColor(0x55, 0x55, 0x88))
    txb(slide, label, x + w/2 - 1.5, y + h/2 - 0.3, 3.0, 0.6,
        size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, italic=True)


# ════════════════════════════════════════════════════
# 슬라이드 1: 표지
# ════════════════════════════════════════════════════
s1 = add_slide(ACCENT1)

# 하단 바
r = s1.shapes.add_shape(1, Inches(0), Inches(7.42), Inches(13.33), Inches(0.08))
r.fill.solid(); r.fill.fore_color.rgb = ACCENT1; r.line.fill.background()

txb(s1, "HandAvatar", 1.0, 1.5, 11.33, 1.5,
    size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s1, "프로젝트 진행 보고",
    1.0, 3.0, 11.33, 0.8, size=30, color=SUBTITLE, align=PP_ALIGN.CENTER)

# 구분선
sep = s1.shapes.add_shape(1, Inches(4.2), Inches(3.95), Inches(4.93), Inches(0.05))
sep.fill.solid(); sep.fill.fore_color.rgb = ACCENT3; sep.line.fill.background()

txb(s1, "손 제스처 → 동물 애니메이션 실시간 제어",
    1.0, 4.2, 11.33, 0.5, size=18, color=LIGHT_GRAY,
    align=PP_ALIGN.CENTER, italic=True)

txb(s1, "2026. 06", 1.0, 6.7, 11.33, 0.5,
    size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# 슬라이드 2: 목차
# ════════════════════════════════════════════════════
s2 = add_slide(ACCENT1)
slide_title(s2, "목차")

items = [
    ("01", "이번 기간 작업 요약",       ACCENT1),
    ("02", "Blend 모드 완성 (5/22 ~ 5/29)",  ACCENT1),
    ("03", "Locomotion 시스템 구현 (6/2 ~ 6/3)", ACCENT2),
    ("04", "손 추적 파이프라인 개선 (6/4)",   ACCENT2),
    ("05", "현재 시스템 시연",           ACCENT3),
    ("06", "미해결 문제",                ACCENT4),
    ("07", "향후 계획",                  LIGHT_GRAY),
]
for i, (num, label, col) in enumerate(items):
    y = 1.4 + i * 0.78
    rect(s2, 0.4, y, 0.7, 0.6, fill=RGBColor(0x1A, 0x1A, 0x30), line=col)
    txb(s2, num, 0.4, y + 0.05, 0.7, 0.5,
        size=17, bold=True, color=col, align=PP_ALIGN.CENTER)
    txb(s2, label, 1.3, y + 0.08, 11.0, 0.5,
        size=20, bold=(i < 4), color=WHITE)


# ════════════════════════════════════════════════════
# 슬라이드 3: 작업 요약 타임라인
# ════════════════════════════════════════════════════
s3 = add_slide(ACCENT2)
slide_title(s3, "이번 기간 작업 요약", "5/22 ~ 6/4 주요 변경 사항", ACCENT2)

rows = [
    ("5 / 22", "Spider 매핑 파이프라인 완전 재구축, KeyframeMappingEngine axis 버그 수정,\nUnity 에디터 UI 개선 (Include 체크박스, 매핑 필터링)", ACCENT1),
    ("5 / 29", "Blend 모드 버그 7개 수정 — restRotation T-포즈, chains 덮어쓰기,\nEuler 180° 경계 점프, Death sin 커브, LDA 임계값 등", ACCENT1),
    ("5 / 30", "시스템 전체 문서화 (blend_mode.md), 발표자료 1차 제작,\nLDA → min_distance 교체 방향 결정", ACCENT3),
    ("6 / 02", "Locomotion 시스템 설계 및 구현 — LocomotionMapper, AnimalLocomotion,\nLocomotionConfigEditor 신규 작성", ACCENT2),
    ("6 / 03", "PythonLauncher (Unity Editor 창), FollowCamera (3인칭) 추가,\n첫 Locomotion 실행 테스트", ACCENT2),
    ("6 / 04", "손 추적 안정화 (palm-local frame, 커플링, DOF delta 트리거),\n이산 속도, velocity deadzone, HUD 추가", ACCENT2),
]

for i, (date, content, col) in enumerate(rows):
    y = 1.35 + i * 1.0
    rect(s3, 0.3, y, 1.4, 0.82,
         fill=RGBColor(0x1A, 0x1A, 0x30), line=col)
    txb(s3, date, 0.3, y + 0.15, 1.4, 0.5,
        size=15, bold=True, color=col, align=PP_ALIGN.CENTER)
    rect(s3, 1.85, y, 11.1, 0.82,
         fill=RGBColor(0x16, 0x16, 0x26), line=RGBColor(0x33, 0x33, 0x55))
    txb(s3, content, 2.0, y + 0.08, 10.8, 0.68, size=14, color=WHITE)


# ════════════════════════════════════════════════════
# 슬라이드 4: Blend 모드 완성
# ════════════════════════════════════════════════════
s4 = add_slide(ACCENT1)
slide_title(s4, "Blend 모드 완성", "주요 버그 7개 수정 — 5/22 ~ 5/29", ACCENT1)

# 왼쪽: 버그 목록
rect(s4, 0.3, 1.3, 7.2, 5.5,
     fill=RGBColor(0x0D, 0x0D, 0x20), line=ACCENT1)
txb(s4, "수정된 버그", 0.5, 1.35, 6.8, 0.45,
    size=17, bold=True, color=ACCENT1)

bugs = [
    ("① restRotation T-포즈",   "Idle→Active 전환 시 재캡처로 수정"),
    ("② body_mapping 분리",     "bone/bodyik_001 매핑 독립, atack1 제외"),
    ("③ Euler 180° 경계 점프",  "최단경로 보간 _short() 추가"),
    ("④ Death sin 커브 복귀",   "후반부 alpha=1.0 유지 (min(t,0.5))"),
    ("⑤ LDA 임계값 오류",       "threshold 739→~16, 트리거 정상화"),
    ("⑥ chains 덮어쓰기",       "depth-tracking 루프로 기존 chains 보존"),
    ("⑦ Direct 엔진 미캘리브",  "_engine_direct.calibrate() 추가"),
]
b4 = s4.shapes.add_textbox(Inches(0.5), Inches(1.85), Inches(6.8), Inches(4.7))
b4.word_wrap = True
tf = b4.text_frame; tf.word_wrap = True
first = True
for title, desc in bugs:
    if first:
        p = tf.paragraphs[0]; first = False
    else:
        p = tf.add_paragraph()
    p.space_before = Pt(6)
    r1 = p.add_run(); r1.text = title + "  "
    r1.font.size = Pt(14); r1.font.bold = True
    r1.font.color.rgb = ACCENT3; r1.font.name = "Malgun Gothic"
    r2 = p.add_run(); r2.text = desc
    r2.font.size = Pt(14); r2.font.color.rgb = WHITE
    r2.font.name = "Malgun Gothic"

# 오른쪽: 상태머신 + 영상
rect(s4, 7.8, 1.3, 5.2, 2.5,
     fill=RGBColor(0x0D, 0x0D, 0x20), line=ACCENT2)
txb(s4, "상태머신", 8.0, 1.35, 4.8, 0.4,
    size=16, bold=True, color=ACCENT2)
txb(s4,
    "  normal\n"
    "  ↓  DOF delta 10f 조건\n"
    "  trigger (Attack1 / Attack2)\n"
    "  ↓  keyframe 재생 완료\n"
    "  cooldown (30f)\n"
    "  ↓\n"
    "  normal",
    8.0, 1.8, 4.8, 1.7, size=13, color=WHITE)

video_placeholder(s4, 7.8, 4.0, 5.2, 2.8,
                  "📹 Attack1/Attack2 시연 영상")


# ════════════════════════════════════════════════════
# 슬라이드 5: Locomotion 시스템
# ════════════════════════════════════════════════════
s5 = add_slide(ACCENT2)
slide_title(s5, "Locomotion 시스템 구현",
            "wrist_dev → 방향 전환 / finger velocity → 이동 속도", ACCENT2)

# 구조 다이어그램 (상단)
steps5 = [
    ("손\n(wrist_dev)",  ACCENT1),
    ("LocomotionMapper\n(Python)",   ACCENT2),
    ("WebSocket\n전송",              ACCENT3),
    ("AnimalLocomotion\n(Unity)",    ACCENT2),
    ("transform\n이동/회전",         ACCENT1),
]
bw = 2.1
for i, (label, col) in enumerate(steps5):
    bx = 0.3 + i * (bw + 0.15)
    rect(s5, bx, 1.3, bw, 1.1,
         fill=RGBColor(0x1A, 0x1A, 0x30), line=col)
    txb(s5, label, bx + 0.05, 1.38, bw - 0.1, 0.95,
        size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    if i < 4:
        txb(s5, "→", bx + bw, 1.65, 0.18, 0.4,
            size=18, bold=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 왼쪽: Python 측
rect(s5, 0.3, 2.65, 6.2, 4.1,
     fill=RGBColor(0x0D, 0x20, 0x0D), line=ACCENT2)
txb(s5, "LocomotionMapper (Python)", 0.5, 2.7, 5.8, 0.45,
    size=16, bold=True, color=ACCENT2)
b5l = s5.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(5.8), Inches(3.3))
b5l.word_wrap = True
tf5l = b5l.text_frame; tf5l.word_wrap = True
lines_l = [
    ("방향 (yaw_delta)", ""),
    ("", "wrist_dev − ref_wrist_dev → deadzone 8°"),
    ("", "deadzone 초과분 × dir_scale → °/frame"),
    ("속도 (speed)", ""),
    ("", "finger velocity (DOF 벡터 L2 노름)"),
    ("", "deadzone 12 이하 노이즈 제거"),
    ("", "EMA 평활화 후 × VEL_SCALE"),
    ("이산화 (Unity)", ""),
    ("", "speed < 0.05 → STOP"),
    ("", "speed ≥ 0.05 → NORMAL (20 units/s)"),
]
first5 = True
for title, desc in lines_l:
    if first5:
        p = tf5l.paragraphs[0]; first5 = False
    else:
        p = tf5l.add_paragraph()
    p.space_before = Pt(3)
    if title and not desc:
        r = p.add_run(); r.text = title
        r.font.size = Pt(14); r.font.bold = True
        r.font.color.rgb = ACCENT3; r.font.name = "Malgun Gothic"
    else:
        r = p.add_run(); r.text = "  " + desc
        r.font.size = Pt(13); r.font.color.rgb = WHITE
        r.font.name = "Malgun Gothic"

# 오른쪽: Unity 측 + 영상
rect(s5, 6.8, 2.65, 6.2, 1.8,
     fill=RGBColor(0x0D, 0x20, 0x0D), line=ACCENT1)
txb(s5, "AnimalLocomotion (Unity)", 7.0, 2.7, 5.8, 0.45,
    size=16, bold=True, color=ACCENT1)
txb(s5,
    "  • STOP / NORMAL SPEED HUD 표시\n"
    "  • TURNING LEFT / RIGHT 표시\n"
    "  • valid=False → 즉시 정지\n"
    "  • FollowCamera 3인칭 추적",
    7.0, 3.2, 5.8, 1.1, size=13, color=WHITE)

video_placeholder(s5, 6.8, 4.6, 6.2, 2.15,
                  "📹 거미 이동/방향전환 시연 영상")


# ════════════════════════════════════════════════════
# 슬라이드 6: 손 추적 개선
# ════════════════════════════════════════════════════
s6 = add_slide(ACCENT2)
slide_title(s6, "손 추적 파이프라인 개선",
            "palm-local 좌표계 · 커플링 · DOF delta 트리거 · 대칭화", ACCENT2)

improvements = [
    ("① palm-local 좌표계",
     "Before: world 3D → 손목 방향에 따라 각도 왜곡\n"
     "After:  검지~소지 = 2D 이미지 좌표 기반\n"
     "        엄지/손목 = world frame 유지\n"
     "→ 손목 회전 무관하게 안정적인 굴신각",
     ACCENT1),
    ("② 손가락 커플링",
     "MCP / PIP / DIP = weighted composite\n"
     "PIP = 1.5×,  DIP = 0.8×\n"
     "→ 손가락 마디 간 자연스러운 연동",
     ACCENT3),
    ("③ DOF delta 트리거",
     "Before: animal joint L2 거리 비교\n"
     "After:  캘리브 기준 DOF 변화량 직접 비교\n"
     "        10프레임 연속 조건 (노이즈 오발동 방지)\n"
     "→ 더 직관적이고 예측 가능한 트리거",
     ACCENT2),
    ("④ 좌우 scale 대칭화",
     "l/r 대응 관절 scale_factor\n"
     "= (left + right) / 2 평균 적용\n"
     "→ 좌우 비대칭 완화",
     ACCENT1),
]

for i, (title, body, col) in enumerate(improvements):
    col_idx = i % 2
    row_idx = i // 2
    x = 0.3 + col_idx * 6.6
    y = 1.3 + row_idx * 2.95
    rect(s6, x, y, 6.35, 2.75,
         fill=RGBColor(0x0D, 0x0D, 0x20), line=col)
    txb(s6, title, x + 0.2, y + 0.1, 5.9, 0.45,
        size=17, bold=True, color=col)
    txb(s6, body, x + 0.2, y + 0.6, 5.9, 1.95,
        size=13, color=WHITE)


# ════════════════════════════════════════════════════
# 슬라이드 7: 시연
# ════════════════════════════════════════════════════
s7 = add_slide(ACCENT3)
slide_title(s7, "현재 시스템 시연",
            "전체 파이프라인 통합 동작", ACCENT3)

video_placeholder(s7, 0.3, 1.3, 12.73, 5.5,
                  "📹 전체 통합 시연 영상\n"
                  "[ 캘리브레이션 → 손가락 → 다리 제어 → 이동/방향 → 공격 트리거 ]")


# ════════════════════════════════════════════════════
# 슬라이드 8: 미해결 문제
# ════════════════════════════════════════════════════
s8 = add_slide(ACCENT4)
slide_title(s8, "미해결 문제", "현재 남아 있는 버그 및 미완성 항목", ACCENT4)

issues = [
    ("🔴",  "Locomotion 속도 불안정",
     "가만히 있어도 가끔 이동 / 멈춤 지연\n"
     "→ MediaPipe 노이즈, velocity deadzone 튜닝 진행 중",
     ACCENT4),
    ("🔴",  "트리거 오발동",
     "가만히 있어도 Attack 발동 / 애니메이션 속도 너무 빠름\n"
     "→ hold=10f로 완화, 근본 원인 추가 분석 필요",
     ACCENT4),
    ("🟡",  "Death 뒤집기 미검증",
     "bone axis Z→X 변경 후 Unity 실제 테스트 미완료\n"
     "→ Play 모드에서 X/Y/Z 수동 확인 예정",
     ACCENT3),
    ("🟡",  "butterfly / fish base_anim 미설정",
     "locomotion_config.json에 해당 동물 base_anim 비어있음",
     ACCENT3),
    ("🟢",  "FollowCamera 거리/높이 튜닝",
     "distanceBack, height Inspector 조정 필요",
     ACCENT2),
]

for i, (icon, title, desc, col) in enumerate(issues):
    y = 1.3 + i * 1.18
    rect(s8, 0.3, y, 0.65, 1.0,
         fill=RGBColor(0x1A, 0x1A, 0x30), line=col)
    txb(s8, icon, 0.3, y + 0.22, 0.65, 0.55,
        size=22, align=PP_ALIGN.CENTER)
    rect(s8, 1.1, y, 11.9, 1.0,
         fill=RGBColor(0x16, 0x16, 0x26), line=col)
    txb(s8, title, 1.25, y + 0.05, 11.5, 0.4,
        size=16, bold=True, color=col)
    txb(s8, desc, 1.25, y + 0.48, 11.5, 0.45,
        size=13, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════
# 슬라이드 9: 향후 계획
# ════════════════════════════════════════════════════
s9 = add_slide(ACCENT1)
slide_title(s9, "향후 계획", "단기 / 중기 목표", ACCENT1)

# 단기
rect(s9, 0.3, 1.3, 6.1, 5.5,
     fill=RGBColor(0x0D, 0x1A, 0x30), line=ACCENT1)
txb(s9, "단기 — 이번 주", 0.5, 1.35, 5.7, 0.45,
    size=18, bold=True, color=ACCENT1)

short_items = [
    "Locomotion 안정화",
    "  • velocity deadzone 최적값 확정",
    "  • idleThreshold 정밀 튜닝",
    "",
    "트리거 안정화",
    "  • hold 프레임 추가 조정",
    "  • 트리거 delta 값 재설정",
    "  • 재생 속도(FPS) 파라미터화",
    "",
    "Death 뒤집기 완성",
    "  • bone axis X Unity 검증",
    "  • 안 되면 Y축 시도",
]
b9l = s9.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(5.7), Inches(4.5))
b9l.word_wrap = True
tf9l = b9l.text_frame; tf9l.word_wrap = True
first9 = True
for line in short_items:
    if first9:
        p = tf9l.paragraphs[0]; first9 = False
    else:
        p = tf9l.add_paragraph()
    p.space_before = Pt(3)
    r = p.add_run()
    r.text = line
    is_header = bool(line) and not line.startswith("  ")
    r.font.size = Pt(15 if is_header else 13)
    r.font.bold = is_header
    r.font.color.rgb = ACCENT3 if is_header else WHITE
    r.font.name = "Malgun Gothic"

# 중기
rect(s9, 6.7, 1.3, 6.3, 5.5,
     fill=RGBColor(0x0D, 0x20, 0x0D), line=ACCENT2)
txb(s9, "중기 — 다음 주 이후", 6.9, 1.35, 5.9, 0.45,
    size=18, bold=True, color=ACCENT2)

mid_items = [
    "말(Horse) 동물 추가",
    "  • locomotion_config horse 설정",
    "  • 매핑 파이프라인 동일 적용",
    "  • 시연 영상 촬영",
    "",
    "전체 시연 준비",
    "  • 안정화된 영상 촬영",
    "  • 다중 동물 전환 시연",
    "",
    "선택 사항",
    "  • LDA → min_distance 교체",
    "  • butterfly / fish 추가",
]
b9r = s9.shapes.add_textbox(Inches(6.9), Inches(1.9), Inches(5.9), Inches(4.5))
b9r.word_wrap = True
tf9r = b9r.text_frame; tf9r.word_wrap = True
first9r = True
for line in mid_items:
    if first9r:
        p = tf9r.paragraphs[0]; first9r = False
    else:
        p = tf9r.add_paragraph()
    p.space_before = Pt(3)
    r = p.add_run()
    r.text = line
    is_header = bool(line) and not line.startswith("  ")
    r.font.size = Pt(15 if is_header else 13)
    r.font.bold = is_header
    r.font.color.rgb = ACCENT3 if is_header else WHITE
    r.font.name = "Malgun Gothic"


# ════════════════════════════════════════════════════
# 슬라이드 10: 마무리
# ════════════════════════════════════════════════════
s10 = add_slide(ACCENT1)
r10 = s10.shapes.add_shape(1, Inches(0), Inches(7.42), Inches(13.33), Inches(0.08))
r10.fill.solid(); r10.fill.fore_color.rgb = ACCENT1; r10.line.fill.background()

txb(s10, "Summary", 1.0, 1.3, 11.33, 1.2,
    size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

sep10 = s10.shapes.add_shape(1, Inches(4.2), Inches(2.65), Inches(4.93), Inches(0.05))
sep10.fill.solid(); sep10.fill.fore_color.rgb = ACCENT3; sep10.line.fill.background()

summary = [
    ("✅ 완료", "Blend 모드 파이프라인 완성 (버그 7개 수정)", ACCENT2),
    ("✅ 완료", "Locomotion 시스템 설계 및 기본 구현",         ACCENT2),
    ("✅ 완료", "손 추적 안정화 (palm-local, 커플링, 이산속도)", ACCENT2),
    ("🔧 진행", "Locomotion / 트리거 튜닝 중",                 ACCENT3),
    ("📌 예정", "Horse 추가, 전체 시연 영상 촬영",              ACCENT1),
]
for i, (icon, text, col) in enumerate(summary):
    y = 3.05 + i * 0.72
    txb(s10, icon, 1.5, y, 1.2, 0.55,
        size=16, bold=True, color=col)
    txb(s10, text, 2.9, y, 9.5, 0.55,
        size=16, color=WHITE)

txb(s10, "HandAnimal Project · 2026",
    1.0, 6.7, 11.33, 0.5, size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ── 저장 ──────────────────────────────────────────
out_path = os.path.normpath(os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    "..", "..", "progress_report.pptx"
))
prs.save(out_path)
print(f"[완료] 저장: {out_path}")
