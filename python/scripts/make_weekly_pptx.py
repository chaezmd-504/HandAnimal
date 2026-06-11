"""
make_weekly_pptx.py
--------------------
주차별 진행 보고 슬라이드 생성 (캔바 옮기기용 — 심플/백지)
실행: python scripts/make_weekly_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

OUT_PATH = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                        "weekly_report.pptx")

W = Inches(13.33)
H = Inches(7.5)

BLACK  = RGBColor(0x1A, 0x1A, 0x1A)
GRAY   = RGBColor(0x88, 0x88, 0x88)
ACCENT = RGBColor(0x2D, 0x6A, 0xF5)  # 파란 포인트 하나만

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]  # 완전 빈 레이아웃
    return prs.slides.add_slide(layout)

def add_text(slide, text, left, top, width, height,
             size=24, bold=False, color=BLACK, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_bullet_box(slide, lines, left, top, width, height, size=18, color=BLACK):
    """lines: list of (indent, text) — indent 0=본문, 1=서브"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (indent, text) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = indent
        run = p.add_run()
        run.text = ("    " * indent) + ("• " if indent == 0 else "- ") + text
        run.font.size = Pt(size if indent == 0 else size - 2)
        run.font.color.rgb = color if indent == 0 else GRAY
    return txBox

def hline(slide, top, color=RGBColor(0xDD,0xDD,0xDD)):
    line = slide.shapes.add_connector(1, Inches(0.6), top, Inches(12.73), top)
    line.line.color.rgb = color
    line.line.width = Pt(0.75)

# ──────────────────────────────────────────────
# Slide 1 — 표지
# ──────────────────────────────────────────────
def slide_cover(prs):
    s = blank_slide(prs)
    add_text(s, "HandAnimal", Inches(0.8), Inches(2.2), Inches(11), Inches(1.5),
             size=52, bold=True, color=BLACK, align=PP_ALIGN.LEFT)
    add_text(s, "주차별 진행 보고  ·  2026-06-05", Inches(0.8), Inches(3.5), Inches(11), Inches(0.8),
             size=22, color=GRAY, align=PP_ALIGN.LEFT)
    hline(s, Inches(3.35), color=ACCENT)

# ──────────────────────────────────────────────
# Slide 2 — Locomotion
# ──────────────────────────────────────────────
def slide_locomotion(prs):
    s = blank_slide(prs)
    add_text(s, "Locomotion", Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             size=36, bold=True, color=BLACK)
    add_text(s, "손으로 동물을 직접 이동시키다", Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
             size=20, color=GRAY, italic=True)
    hline(s, Inches(1.55))

    bullets = [
        (0, "방향 제어 — 오른손 손목 좌우 기울기 → yaw 회전"),
        (0, "속도 계산 — 손 포즈가 Walk 기준 포즈에 가까울수록 speed 증가"),
        (0, "Head Direction — 얼굴이 향하는 방향으로 이동 방향 제어"),
        (0, "Walk 애니메이션 — 이동 중 자동 재생 / 정지 시 멈춤"),
    ]
    add_bullet_box(s, bullets, Inches(0.8), Inches(1.8), Inches(11.5), Inches(4), size=20)

    add_text(s, "[ 다이어그램: 손목→yaw / Walk 포즈→speed 흐름도 ]",
             Inches(0.8), Inches(6.3), Inches(11), Inches(0.6),
             size=14, color=GRAY, italic=True)

# ──────────────────────────────────────────────
# Slide 3 — Trigger 시스템
# ──────────────────────────────────────────────
def slide_trigger(prs):
    s = blank_slide(prs)
    add_text(s, "Trigger 시스템", Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             size=36, bold=True, color=BLACK)
    add_text(s, "특정 동작으로 애니메이션 발동", Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
             size=20, color=GRAY, italic=True)
    hline(s, Inches(1.55))

    bullets = [
        (0, "Snap — 손가락이 빠르게 구부러지는 순간 즉시 발동 (확! 동작)"),
        (0, "Hold — 일정 각도 이상 3프레임 유지 시 발동"),
        (0, "Unity Animator 직접 연동 → 실제 클립 재생 (keyframe 방식 대체)"),
        (0, "재생 중 손 입력도 blend 비율로 반영"),
    ]
    add_bullet_box(s, bullets, Inches(0.8), Inches(1.8), Inches(11.5), Inches(4), size=20)

    add_text(s, "[ 영상: Attack1/2 발동 장면 ]",
             Inches(0.8), Inches(6.3), Inches(11), Inches(0.6),
             size=14, color=GRAY, italic=True)

# ──────────────────────────────────────────────
# Slide 4 — 현재 상태 데모
# ──────────────────────────────────────────────
def slide_demo(prs):
    s = blank_slide(prs)
    add_text(s, "현재 상태 데모", Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             size=36, bold=True, color=BLACK)
    add_text(s, "Spider 통합 동작", Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
             size=20, color=GRAY, italic=True)
    hline(s, Inches(1.55))

    bullets = [
        (0, "손 감지 → 다리 매핑 → Walk 이동 → 공격 트리거"),
        (0, "얼굴 방향으로 이동 방향 전환"),
    ]
    add_bullet_box(s, bullets, Inches(0.8), Inches(1.8), Inches(11.5), Inches(1.8), size=20)

    # 영상 자리 표시 박스
    box = s.shapes.add_shape(1,  # MSO_SHAPE_TYPE.RECTANGLE
                              Inches(1.5), Inches(3.0), Inches(10.3), Inches(3.8))
    box.fill.background()
    box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    box.line.width = Pt(1)
    txBox = box.text_frame
    txBox.text = "[ 영상 삽입 ]"
    p = txBox.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(18)
    run.font.color.rgb = GRAY
    run.font.italic = True

# ──────────────────────────────────────────────
# Slide 5 — 다음 단계
# ──────────────────────────────────────────────
def slide_next(prs):
    s = blank_slide(prs)
    add_text(s, "다음 단계", Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             size=36, bold=True, color=BLACK)
    hline(s, Inches(1.15))

    bullets = [
        (0, "Horse 세팅 및 테스트"),
        (1, "왼손 4손가락 = 앞뒤 4다리"),
        (1, "오른손 손목 = 머리 / 목"),
        (0, "매핑 품질 개선"),
        (1, "알고리즘 손가락 분리 제약 추가"),
    ]
    add_bullet_box(s, bullets, Inches(0.8), Inches(1.5), Inches(11.5), Inches(4.5), size=22)

# ──────────────────────────────────────────────
# 생성
# ──────────────────────────────────────────────
def main():
    prs = new_prs()
    slide_cover(prs)
    slide_locomotion(prs)
    slide_trigger(prs)
    slide_demo(prs)
    slide_next(prs)
    prs.save(OUT_PATH)
    print(f"저장됨: {OUT_PATH}")

if __name__ == "__main__":
    main()
