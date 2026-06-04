"""
make_blend_pptx.py
Blend 모드 발표용 PowerPoint 자동 생성
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── 색상 팔레트 ──────────────────────────────────────
BG_DARK    = RGBColor(0x12, 0x12, 0x1E)   # 거의 검정 (배경)
ACCENT1    = RGBColor(0x5B, 0x9B, 0xD5)   # 파랑 (Direct 엔진)
ACCENT2    = RGBColor(0x70, 0xAD, 0x47)   # 초록 (Sequential 엔진)
ACCENT3    = RGBColor(0xFF, 0xC0, 0x00)   # 노랑 (Alpha / 트리거)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_GRAY  = RGBColor(0x2D, 0x2D, 0x3F)
SUBTITLE   = RGBColor(0xA0, 0xC4, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # 완전 빈 레이아웃


def add_slide():
    slide = prs.slides.add_slide(BLANK)
    # 배경 단색 채우기
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK
    return slide


def txb(slide, text, x, y, w, h,
        size=24, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Malgun Gothic"
    return box


def rect(slide, x, y, w, h, fill_color=DARK_GRAY,
         line_color=None, line_width=Pt(1.5), radius=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_paragraph(txbox, text, size=18, bold=False, color=WHITE,
                  align=PP_ALIGN.LEFT, italic=False, space_before=Pt(4)):
    tf = txbox.text_frame
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Malgun Gothic"
    return p


# ════════════════════════════════════════════════════
# 슬라이드 1: 표지
# ════════════════════════════════════════════════════
s1 = add_slide()

# 상단 장식 바
rect(s1, 0, 0, 13.33, 0.08, fill_color=ACCENT1)
rect(s1, 0, 7.42, 13.33, 0.08, fill_color=ACCENT1)

# 제목
txb(s1, "Blend Mode", 1.5, 1.8, 10, 1.4,
    size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 부제목
txb(s1, "손 제스처 → 동물 애니메이션 블렌딩 파이프라인",
    1.5, 3.3, 10, 0.8, size=26, color=SUBTITLE, align=PP_ALIGN.CENTER)

# 구분선
r = s1.shapes.add_shape(1, Inches(4.5), Inches(4.15), Inches(4.33), Inches(0.04))
r.fill.solid(); r.fill.fore_color.rgb = ACCENT3
r.line.fill.background()

# 설명
txb(s1, "Direct Engine  +  Sequential Engine  +  Alpha Blending",
    1.5, 4.4, 10, 0.6, size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, italic=True)

txb(s1, "HandAnimal Project · 2026",
    1.5, 6.7, 10, 0.5, size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# 슬라이드 2: Blend 모드란?
# ════════════════════════════════════════════════════
s2 = add_slide()
rect(s2, 0, 0, 13.33, 0.08, fill_color=ACCENT1)

txb(s2, "Blend 모드란?", 0.4, 0.2, 10, 0.7,
    size=36, bold=True, color=WHITE)
txb(s2, "두 가지 요구사항을 동시에 만족시키는 하이브리드 방식",
    0.4, 0.85, 12, 0.5, size=18, color=SUBTITLE)

# 왼쪽 박스: 연속성
rect(s2, 0.4, 1.6, 5.8, 2.2, fill_color=RGBColor(0x0D, 0x2A, 0x45),
     line_color=ACCENT1)
txb(s2, "연속성 (Continuous)", 0.6, 1.65, 5.4, 0.5,
    size=20, bold=True, color=ACCENT1)
txb(s2, "손가락 굽힘 → 다리가 즉시 반응\n손 DOF × scale → 관절 각도\n(반응성, 실시간)",
    0.6, 2.15, 5.4, 1.4, size=17, color=WHITE)

# 오른쪽 박스: 특수 동작
rect(s2, 7.1, 1.6, 5.8, 2.2, fill_color=RGBColor(0x0D, 0x2D, 0x12),
     line_color=ACCENT2)
txb(s2, "특수 동작 (Action)", 7.3, 1.65, 5.4, 0.5,
    size=20, bold=True, color=ACCENT2)
txb(s2, "주먹 쥐기 / 손목 비틀기 →\n사전 디자인된 애니메이션 재생\n(자연스러운 액션)",
    7.3, 2.15, 5.4, 1.4, size=17, color=WHITE)

# 가운데 플러스
txb(s2, "+", 6.3, 2.1, 0.7, 1.0,
    size=48, bold=True, color=ACCENT3, align=PP_ALIGN.CENTER)

# 화살표 → Blend
txb(s2, "↓", 6.3, 3.6, 0.7, 0.6,
    size=36, bold=True, color=ACCENT3, align=PP_ALIGN.CENTER)

# Blend 결과 박스
rect(s2, 3.0, 4.3, 7.3, 1.2, fill_color=RGBColor(0x2D, 0x1E, 0x00),
     line_color=ACCENT3)
txb(s2, "Blend =  (1 - α) × joints_direct  +  α × joints_keyframe",
    3.2, 4.55, 7.0, 0.6, size=20, bold=True, color=ACCENT3, align=PP_ALIGN.CENTER)

# 하단 설명
txb(s2, "α = 0 → 손 제어 100%     α = 1 → 애니메이션 100%     0<α<1 → 자연스러운 혼합",
    0.4, 5.7, 12.5, 0.5, size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# 슬라이드 3: 전체 파이프라인 흐름
# ════════════════════════════════════════════════════
s3 = add_slide()
rect(s3, 0, 0, 13.33, 0.08, fill_color=ACCENT2)

txb(s3, "전체 파이프라인 흐름", 0.4, 0.2, 10, 0.7,
    size=36, bold=True, color=WHITE)

# 흐름 박스들 (좌→우 5단계)
steps = [
    ("웹캠\n프레임", ACCENT1),
    ("MediaPipe\n손 랜드마크\n(21점×2손)", ACCENT1),
    ("DOF 계산\n+EMA 스무딩\n(20 DOF/손)", ACCENT1),
    ("Direct\n엔진\n(즉시 변환)", ACCENT1),
    ("LDA\n트리거\n판별", ACCENT3),
]
box_w = 2.0
box_h = 1.6
start_x = 0.35
gap = 0.22

for i, (label, col) in enumerate(steps):
    bx = start_x + i * (box_w + gap)
    rect(s3, bx, 1.4, box_w, box_h,
         fill_color=RGBColor(0x1A, 0x1A, 0x30), line_color=col)
    txb(s3, label, bx + 0.05, 1.5, box_w - 0.1, box_h - 0.1,
        size=15, bold=True, color=col, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        txb(s3, "→", bx + box_w, 1.85, gap + 0.05, 0.7,
            size=22, bold=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 두 갈래 분기
txb(s3, "↓  트리거 없음", 1.3, 3.1, 3.5, 0.4, size=15, color=ACCENT1)
txb(s3, "↓  트리거 발동", 7.8, 3.1, 3.5, 0.4, size=15, color=ACCENT3)

# 왼쪽 (continuous)
rect(s3, 0.35, 3.55, 5.5, 1.0, fill_color=RGBColor(0x0A, 0x20, 0x3A), line_color=ACCENT1)
txb(s3, "joints = joints_direct   (α = 0)", 0.5, 3.7, 5.2, 0.7,
    size=17, bold=True, color=ACCENT1, align=PP_ALIGN.CENTER)

# 오른쪽 (action)
rect(s3, 7.48, 3.55, 5.5, 1.0, fill_color=RGBColor(0x1A, 0x2A, 0x00), line_color=ACCENT2)
txb(s3, "Sequential 키프레임 보간\n→ Alpha 블렌딩 적용", 7.6, 3.6, 5.2, 0.85,
    size=15, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

# 합류 화살표
txb(s3, "↘", 4.5, 4.55, 1.5, 0.5, size=28, color=ACCENT3, align=PP_ALIGN.CENTER)
txb(s3, "↙", 7.4, 4.55, 1.5, 0.5, size=28, color=ACCENT3, align=PP_ALIGN.CENTER)

# WebSocket 전송
rect(s3, 3.9, 5.15, 5.5, 0.85, fill_color=RGBColor(0x2D, 0x1A, 0x3A), line_color=ACCENT3)
txb(s3, "WebSocket 전송 → Unity AnimalController",
    4.1, 5.3, 5.1, 0.55, size=16, bold=True, color=ACCENT3, align=PP_ALIGN.CENTER)

txb(s3, '{"l_leg":{"x":0,"y":0,"z":-32.5}, "r_leg":{...}, ...}',
    0.4, 6.2, 12.5, 0.5, size=13, italic=True,
    color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# 슬라이드 4: 두 엔진 비교
# ════════════════════════════════════════════════════
s4 = add_slide()
rect(s4, 0, 0, 13.33, 0.08, fill_color=ACCENT1)

txb(s4, "두 엔진: Direct vs Sequential", 0.4, 0.2, 12, 0.7,
    size=36, bold=True, color=WHITE)

# Direct 엔진 박스
rect(s4, 0.3, 1.1, 6.0, 5.6, fill_color=RGBColor(0x0A, 0x1F, 0x38), line_color=ACCENT1)
txb(s4, "Direct Engine", 0.5, 1.15, 5.6, 0.55,
    size=22, bold=True, color=ACCENT1)
txb(s4, "(MappingEngine)", 0.5, 1.65, 5.6, 0.4,
    size=14, italic=True, color=LIGHT_GRAY)

b4l = slide_txb = s4.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(5.6), Inches(4.3))
b4l.word_wrap = True
tf = b4l.text_frame; tf.word_wrap = True
items_left = [
    ("역할: ", "손 DOF를 즉각 관절 각도로 변환", False),
    ("", "", False),
    ("수식:", "", True),
    ("joint = (dof - dof_ref) × scale", "", False),
    ("", "", False),
    ("예시 (l_leg):", "", True),
    ("index_mcp = 45°,  ref = 12°", "", False),
    ("scale = 1.95", "", False),
    ("→  l_leg = (45-12) × 1.95 = 64.35°", "", False),
    ("", "", False),
    ("특징:", "", True),
    ("매 프레임 즉시 계산 (반응성 최고)", "", False),
    ("노이즈 → 3° 데드존으로 보완", "", False),
]
first = True
for label, val, is_header in items_left:
    if first:
        p = tf.paragraphs[0]; first = False
    else:
        p = tf.add_paragraph()
    p.space_before = Pt(2)
    if label and val:
        r1 = p.add_run(); r1.text = label
        r1.font.size = Pt(15); r1.font.bold = True
        r1.font.color.rgb = ACCENT1; r1.font.name = "Malgun Gothic"
        r2 = p.add_run(); r2.text = val
        r2.font.size = Pt(15); r2.font.color.rgb = WHITE; r2.font.name = "Malgun Gothic"
    else:
        r = p.add_run()
        r.text = label + val
        r.font.size = Pt(15 if not is_header else 16)
        r.font.bold = is_header
        r.font.color.rgb = ACCENT1 if is_header else WHITE
        r.font.name = "Malgun Gothic"

# Sequential 엔진 박스
rect(s4, 7.0, 1.1, 6.0, 5.6, fill_color=RGBColor(0x0A, 0x28, 0x0A), line_color=ACCENT2)
txb(s4, "Sequential Engine", 7.2, 1.15, 5.6, 0.55,
    size=22, bold=True, color=ACCENT2)
txb(s4, "(KeyframeMappingEngine)", 7.2, 1.65, 5.6, 0.4,
    size=14, italic=True, color=LIGHT_GRAY)

b4r = s4.shapes.add_textbox(Inches(7.2), Inches(2.1), Inches(5.6), Inches(4.3))
b4r.word_wrap = True
tf2 = b4r.text_frame; tf2.word_wrap = True
items_right = [
    ("역할: ", "키프레임 애니메이션을 cursor 위치에서 보간"),
    ("", ""),
    ("구조:", ""),
    ("", "[frame0]→[frame1]→...→[frameN]"),
    ("", "cursor 0.0 ~ N-1 (실수)"),
    ("", ""),
    ("보간 수식:", ""),
    ("", "lo = floor(cursor)"),
    ("", "t  = cursor - lo"),
    ("", "pose = lerp(frame[lo], frame[hi], t)"),
    ("", ""),
    ("특징:", ""),
    ("", "Euler 최단 경로 보간 (_short)"),
    ("", "±180° 경계 점프 방지"),
]
first2 = True
for label, val in items_right:
    is_header = label.endswith(":")
    if first2:
        p = tf2.paragraphs[0]; first2 = False
    else:
        p = tf2.add_paragraph()
    p.space_before = Pt(2)
    if label and val:
        r1 = p.add_run(); r1.text = label
        r1.font.size = Pt(15); r1.font.bold = True
        r1.font.color.rgb = ACCENT2; r1.font.name = "Malgun Gothic"
        r2 = p.add_run(); r2.text = val
        r2.font.size = Pt(15); r2.font.color.rgb = WHITE; r2.font.name = "Malgun Gothic"
    else:
        r = p.add_run()
        r.text = label + val
        r.font.size = Pt(15)
        r.font.bold = is_header
        r.font.color.rgb = ACCENT2 if is_header else WHITE
        r.font.name = "Malgun Gothic"


# ════════════════════════════════════════════════════
# 슬라이드 5: LDA 트리거 판별
# ════════════════════════════════════════════════════
s5 = add_slide()
rect(s5, 0, 0, 13.33, 0.08, fill_color=ACCENT3)

txb(s5, "LDA 트리거 판별", 0.4, 0.2, 10, 0.7,
    size=36, bold=True, color=WHITE)
txb(s5, "현재 손 포즈가 '공격 제스처'인가를 수학적으로 판별",
    0.4, 0.85, 12, 0.45, size=17, color=SUBTITLE)

# Step 박스들
steps5 = [
    ("① 손 변화량 계산", "Δh = h_current − h_reference\n(캘리브레이션 기준 포즈와의 차이)", ACCENT1),
    ("② LDA 스코어", "score = w_left · Δh_left[idx]\n        + w_right · Δh_right[idx]", ACCENT2),
    ("③ 임계값 비교", "threshold = separation × 0.5 × k\nscore > threshold  →  트리거!", ACCENT3),
]
bw5 = 3.9
for i, (title, body, col) in enumerate(steps5):
    bx5 = 0.3 + i * (bw5 + 0.22)
    rect(s5, bx5, 1.5, bw5, 2.2,
         fill_color=RGBColor(0x18, 0x18, 0x28), line_color=col)
    txb(s5, title, bx5+0.1, 1.55, bw5-0.2, 0.5,
        size=17, bold=True, color=col)
    txb(s5, body, bx5+0.1, 2.1, bw5-0.2, 1.4, size=15, color=WHITE)
    if i < 2:
        txb(s5, "→", bx5+bw5, 2.35, 0.25, 0.5,
            size=24, bold=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 직관적 해석
rect(s5, 0.3, 3.9, 12.73, 1.1, fill_color=RGBColor(0x20, 0x18, 0x00), line_color=ACCENT3)
txb(s5, "직관적 의미", 0.5, 3.95, 3, 0.4, size=16, bold=True, color=ACCENT3)
txb(s5, '"idle 포즈와 action 포즈를 가장 잘 구별하는 방향으로\n 손 변화량을 투영했을 때의 크기"',
    0.5, 4.3, 12.3, 0.6, size=15, italic=True, color=LIGHT_GRAY)

# 수직선 그림
txb(s5, "스코어 분포:", 0.3, 5.15, 2.5, 0.4, size=15, bold=True, color=WHITE)
txb(s5,
    "      0                  threshold               separation",
    0.3, 5.5, 12.5, 0.35, size=13, color=LIGHT_GRAY)
txb(s5,
    "──────────────── idle ──────────┤──────── action ────────────",
    0.3, 5.8, 12.5, 0.35, size=13, color=LIGHT_GRAY)
txb(s5,
    "k = ATTACK_THRESHOLD (1.0=기본 / 0.7=예민 / 1.5=둔감)",
    0.3, 6.2, 12.5, 0.4, size=14, italic=True, color=ACCENT3)

# wrist_rot 보조
rect(s5, 0.3, 6.65, 12.73, 0.65, fill_color=RGBColor(0x18, 0x18, 0x28), line_color=LIGHT_GRAY)
txb(s5, "wrist_rot 직접 체크 (Death 옵션):  |wrist_rot − ref| ≥ threshold  →  Death 트리거",
    0.5, 6.7, 12.3, 0.5, size=14, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════
# 슬라이드 6: Alpha 커브 & 블렌딩 수식
# ════════════════════════════════════════════════════
s6 = add_slide()
rect(s6, 0, 0, 13.33, 0.08, fill_color=ACCENT2)

txb(s6, "Alpha 커브 & 블렌딩 수식", 0.4, 0.2, 12, 0.7,
    size=36, bold=True, color=WHITE)

# 왼쪽: Attack alpha
rect(s6, 0.3, 1.1, 5.9, 4.0, fill_color=RGBColor(0x0A, 0x1F, 0x38), line_color=ACCENT1)
txb(s6, "Attack (Attack1 / Attack2)", 0.5, 1.15, 5.5, 0.45,
    size=18, bold=True, color=ACCENT1)
txb(s6, "α = sin( cursor/total × π )",
    0.5, 1.65, 5.5, 0.5, size=17, bold=True, color=ACCENT3)
txb(s6,
    "α\n1.0 │     ╭───╮\n     │    ╱     ╲\n0.5 │   ╱       ╲\n     │  ╱         ╲\n0.0 ┼──╱─────────────╲──",
    0.5, 2.2, 5.5, 2.3, size=14, color=WHITE)
txb(s6, "         0              1   (cursor/total)",
    0.5, 4.35, 5.5, 0.4, size=12, color=LIGHT_GRAY)
txb(s6, "0→1→0  : 자연스럽게 진입 후 손 제어로 복귀",
    0.5, 4.7, 5.5, 0.3, size=13, italic=True, color=LIGHT_GRAY)

# 오른쪽: Death alpha
rect(s6, 7.13, 1.1, 5.9, 4.0, fill_color=RGBColor(0x0A, 0x28, 0x0A), line_color=ACCENT2)
txb(s6, "Death (특수 처리)", 7.33, 1.15, 5.5, 0.45,
    size=18, bold=True, color=ACCENT2)
txb(s6, "α = sin( min(cursor/total, 0.5) × π )",
    7.33, 1.65, 5.5, 0.5, size=17, bold=True, color=ACCENT3)
txb(s6,
    "α\n1.0 │     ╭──────────\n     │    ╱\n0.5 │   ╱\n     │  ╱\n0.0 ┼──╱────────────────",
    7.33, 2.2, 5.5, 2.3, size=14, color=WHITE)
txb(s6, "         0              1   (cursor/total)",
    7.33, 4.35, 5.5, 0.4, size=12, color=LIGHT_GRAY)
txb(s6, "0→1→유지  : 한번 발동하면 되돌아오지 않음",
    7.33, 4.7, 5.5, 0.3, size=13, italic=True, color=LIGHT_GRAY)

# 블렌딩 수식 하단
rect(s6, 0.3, 5.3, 12.73, 1.75, fill_color=RGBColor(0x20, 0x18, 0x00), line_color=ACCENT3)
txb(s6, "최종 블렌딩 수식", 0.5, 5.35, 4, 0.45, size=17, bold=True, color=ACCENT3)
txb(s6,
    "joints_final[jid] = (1 - α) × joints_direct[jid]  +  α × joints_keyframe[jid]",
    0.5, 5.8, 12.3, 0.5, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s6,
    "α = 0  손 제어 100%     α = 0.5  50:50 혼합     α = 1  애니메이션 100%",
    0.5, 6.3, 12.3, 0.45, size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# 슬라이드 7: 상태머신
# ════════════════════════════════════════════════════
s7 = add_slide()
rect(s7, 0, 0, 13.33, 0.08, fill_color=ACCENT3)

txb(s7, "상태머신 (State Machine)", 0.4, 0.2, 10, 0.7,
    size=36, bold=True, color=WHITE)

# CONTINUOUS 박스
rect(s7, 0.5, 1.2, 3.8, 1.6, fill_color=RGBColor(0x0A, 0x1F, 0x38), line_color=ACCENT1)
txb(s7, "CONTINUOUS", 0.7, 1.25, 3.4, 0.5, size=20, bold=True, color=ACCENT1)
txb(s7, "joints = joints_direct\nα = 0  (손 제어 100%)", 0.7, 1.7, 3.4, 0.85,
    size=14, color=WHITE)

# 화살표 →
txb(s7, "→", 4.4, 1.8, 0.8, 0.5, size=28, bold=True, color=ACCENT3, align=PP_ALIGN.CENTER)
txb(s7, "LDA score > threshold\n또는 |Δwrist_rot| ≥ thr",
    4.25, 1.35, 1.5, 0.7, size=10, color=ACCENT3, align=PP_ALIGN.CENTER)

# ACTION 박스
rect(s7, 5.3, 1.2, 3.8, 1.6, fill_color=RGBColor(0x1A, 0x2A, 0x00), line_color=ACCENT2)
txb(s7, "ACTION", 5.5, 1.25, 3.4, 0.5, size=20, bold=True, color=ACCENT2)
txb(s7, "cursor 0 → N\nα = sin curve\nblend 적용", 5.5, 1.7, 3.4, 0.85,
    size=14, color=WHITE)

# 화살표 →
txb(s7, "→", 9.2, 1.8, 0.8, 0.5, size=28, bold=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
txb(s7, "cursor ≥ N\n(완료)", 9.1, 1.35, 1.0, 0.7, size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# COOLDOWN 박스
rect(s7, 10.1, 1.2, 2.8, 1.6, fill_color=RGBColor(0x28, 0x1A, 0x00), line_color=LIGHT_GRAY)
txb(s7, "COOLDOWN", 10.3, 1.25, 2.4, 0.5, size=18, bold=True, color=LIGHT_GRAY)
txb(s7, "45 frames\n(~1.5초)", 10.3, 1.7, 2.4, 0.85, size=14, color=WHITE)

# 화살표 되돌아오기 (아래로)
txb(s7, "↓ 종료 후", 11.0, 2.85, 1.5, 0.4, size=11, color=LIGHT_GRAY)
txb(s7, "──────────────────────────────────────────────────────────────────",
    0.5, 3.3, 12.0, 0.3, size=11, color=RGBColor(0x44, 0x44, 0x55))
txb(s7, "↑  cooldown == 0  →  CONTINUOUS 복귀",
    3.0, 3.5, 7.0, 0.4, size=14, bold=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 상세 설명 박스들
details = [
    ("CONTINUOUS 상세", ACCENT1,
     "• 매 프레임 Direct 엔진으로 관절 계산\n• LDA 스코어 감시 중\n• cooldown > 0이면 트리거 무시"),
    ("ACTION 상세", ACCENT2,
     "• cursor += cursor_speed (기본 1.0/frame)\n• get_sequential_pose(anim, cursor) 호출\n• (1-α)×direct + α×keyframe 혼합"),
    ("COOLDOWN 상세", LIGHT_GRAY,
     "• joints = joints_direct (손 제어)\n• 45프레임 대기 후 CONTINUOUS 복귀\n• 연속 재발동 방지"),
]
for i, (title, col, body) in enumerate(details):
    bx7 = 0.3 + i * 4.45
    rect(s7, bx7, 4.0, 4.2, 2.9, fill_color=RGBColor(0x18, 0x18, 0x28), line_color=col)
    txb(s7, title, bx7+0.15, 4.05, 3.9, 0.45, size=15, bold=True, color=col)
    txb(s7, body, bx7+0.15, 4.55, 3.9, 2.1, size=13, color=WHITE)


# ════════════════════════════════════════════════════
# 슬라이드 8: 캘리브레이션
# ════════════════════════════════════════════════════
s8 = add_slide()
rect(s8, 0, 0, 13.33, 0.08, fill_color=ACCENT1)

txb(s8, "캘리브레이션 — 양쪽 엔진 동기화", 0.4, 0.2, 12, 0.7,
    size=34, bold=True, color=WHITE)
txb(s8, "사용자마다 다른 손 크기와 자연 자세를 기준으로 보정", 0.4, 0.85, 12, 0.4,
    size=16, color=SUBTITLE)

# 왜 필요한가
rect(s8, 0.3, 1.4, 4.0, 1.8, fill_color=RGBColor(0x18, 0x18, 0x28), line_color=ACCENT3)
txb(s8, "왜 필요한가?", 0.5, 1.45, 3.6, 0.45, size=16, bold=True, color=ACCENT3)
txb(s8, "• 손 크기/가동범위 개인차\n• 자연 이완 자세 ≠ 완전 편 자세\n• 기준 없이는 캘리브 포즈에서\n  다리가 기울어짐",
    0.5, 1.9, 3.6, 1.1, size=13, color=WHITE)

# 수식
rect(s8, 4.6, 1.4, 8.4, 1.8, fill_color=RGBColor(0x18, 0x18, 0x28), line_color=ACCENT1)
txb(s8, "수식 변화", 4.8, 1.45, 3.6, 0.45, size=16, bold=True, color=ACCENT1)
txb(s8, "캘리브 전:  joint = (dof − h_ref_original) × scale\n캘리브 후:  joint = (dof − h_ref_new) × scale\n\nh_ref_new = 캘리브레이션 시점 실제 DOF 값\n→ 캘리브 포즈에서 joint = 0 보장",
    4.8, 1.9, 7.8, 1.2, size=13, color=WHITE)

# 버그 설명
rect(s8, 0.3, 3.4, 12.73, 1.7, fill_color=RGBColor(0x38, 0x0A, 0x0A), line_color=RGBColor(0xFF, 0x44, 0x44))
txb(s8, "수정된 버그: Direct 엔진 미캘리브레이션", 0.5, 3.45, 8, 0.45,
    size=16, bold=True, color=RGBColor(0xFF, 0x66, 0x66))
txb(s8,
    "이전:  engine.calibrate()만 호출  →  _engine_direct는 원본 h_ref 그대로\n"
    "결과:  캘리브 포즈에서 joints_direct ≠ 0  →  idle 상태인데 다리가 기울어지는 버그",
    0.5, 3.9, 12.3, 0.85, size=13, color=RGBColor(0xFF, 0xCC, 0xCC))

# 수정 코드
rect(s8, 0.3, 5.25, 12.73, 1.7, fill_color=RGBColor(0x0A, 0x20, 0x10), line_color=ACCENT2)
txb(s8, "수정 코드", 0.5, 5.3, 3, 0.4, size=15, bold=True, color=ACCENT2)
txb(s8,
    "engine.calibrate(hands_angles)           # KeyframeMappingEngine\n"
    "if args.mapping == 'blend':\n"
    "    _engine_direct.calibrate(hands_angles)   # MappingEngine  ← 추가",
    0.5, 5.7, 12.3, 1.0, size=14, color=RGBColor(0xA8, 0xFF, 0xC0))


# ════════════════════════════════════════════════════
# 슬라이드 9: 핵심 공식 3가지 요약
# ════════════════════════════════════════════════════
s9 = add_slide()
rect(s9, 0, 0, 13.33, 0.08, fill_color=ACCENT3)

txb(s9, "핵심 공식 3가지", 0.4, 0.2, 10, 0.7,
    size=40, bold=True, color=WHITE)
txb(s9, "Blend 모드의 모든 수식을 한 페이지로",
    0.4, 0.85, 12, 0.4, size=17, color=SUBTITLE)

formulas = [
    ("① Direct 변환", ACCENT1,
     "joint[jid]  =  (hand_dof[assigned_dof]  −  h_ref[assigned_dof])  ×  scale_factor",
     "손 DOF의 캘리브레이션 기준 대비 변화량을 scale로 증폭 → 관절 각도"),
    ("② LDA 트리거", ACCENT3,
     "score = w_left · Δh_left[idx]  +  w_right · Δh_right[idx]\ntrigger if  score  >  separation × 0.5 × k",
     "손 변화량 벡터를 LDA 가중치로 투영 → idle/action 구분 선 초과 시 발동"),
    ("③ Alpha 블렌딩", ACCENT2,
     "joints_final  =  (1 − α) × joints_direct  +  α × joints_keyframe\nα = sin(cursor/total × π)  또는  sin(min(t, 0.5) × π)",
     "두 엔진의 결과를 alpha로 선형 혼합 → 자연스러운 진입/복귀"),
]

for i, (title, col, formula, desc) in enumerate(formulas):
    y = 1.5 + i * 1.9
    rect(s9, 0.3, y, 12.73, 1.7, fill_color=RGBColor(0x18, 0x18, 0x28), line_color=col)
    txb(s9, title, 0.5, y+0.05, 3, 0.45, size=18, bold=True, color=col)
    txb(s9, formula, 0.5, y+0.5, 12.3, 0.65, size=16, bold=True, color=WHITE)
    txb(s9, desc, 0.5, y+1.2, 12.3, 0.38, size=12, italic=True, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════
# 슬라이드 10: 마무리
# ════════════════════════════════════════════════════
s10 = add_slide()
rect(s10, 0, 0, 13.33, 0.08, fill_color=ACCENT1)
rect(s10, 0, 7.42, 13.33, 0.08, fill_color=ACCENT1)

txb(s10, "Blend Mode Summary", 1.5, 1.5, 10, 1.2,
    size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

r = s10.shapes.add_shape(1, Inches(4.5), Inches(2.85), Inches(4.33), Inches(0.04))
r.fill.solid(); r.fill.fore_color.rgb = ACCENT3; r.line.fill.background()

points = [
    ("Direct Engine", "매 프레임 손 DOF → 관절 즉시 변환 (반응성)", ACCENT1),
    ("Sequential Engine", "cursor 기반 키프레임 선형 보간 (부드러운 액션)", ACCENT2),
    ("LDA Trigger", "손 변화량 투영 스코어로 action 자동 감지", ACCENT3),
    ("Alpha Blending", "두 엔진을 sin 커브로 부드럽게 혼합", WHITE),
    ("Calibration", "양쪽 엔진 동시 보정 → idle = 기본 자세 보장", SUBTITLE),
]

for i, (key, val, col) in enumerate(points):
    y = 3.15 + i * 0.67
    txb(s10, f"■  {key}:", 1.8, y, 3.2, 0.5, size=15, bold=True, color=col)
    txb(s10, val, 5.0, y, 7.5, 0.5, size=15, color=WHITE)

txb(s10, "HandAnimal Project · 2026",
    1.5, 6.7, 10, 0.5, size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ── 저장 ──────────────────────────────────────────
import os
out_path = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
    "..", "blend_mode_presentation.pptx"
)
out_path = os.path.normpath(out_path)
prs.save(out_path)
print(f"[완료] 저장: {out_path}")
